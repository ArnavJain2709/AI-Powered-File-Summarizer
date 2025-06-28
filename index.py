import os
import streamlit as st
import fitz  # PyMuPDF
import pptx
import openpyxl
import docx
import requests
import re
import logging

# Configure logging to save to file AND console
logging.basicConfig(
    level=logging.DEBUG, 
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('app.log'),  # Save to app.log file
        logging.StreamHandler()  # Also show in console
    ]
)
logger = logging.getLogger(__name__)

# --- Example Usage and Logging Demo ---
def show_example():
    """Display an example of how to use the application and demonstrate logging."""
    st.header("📚 Example Usage")
    
    with st.expander("🔍 How to Use This Application", expanded=False):
        st.markdown("""
        ### Step-by-Step Guide:
        
        1. **Get Your API Key**: 
           - Visit [Google AI Studio](https://aistudio.google.com/app/apikey)
           - Create a free Gemini API key
           
        2. **Enter Your API Key**: 
           - Paste it in the sidebar (it will be hidden for security)
           
        3. **Choose a Directory**: 
           - Enter a path like: `C:\\Users\\YourName\\Documents`
           - Or: `C:\\Arnav\\Projects` (any folder with files)
           
        4. **Scan & Summarize**: 
           - Click the "Scan & Summarize Directory" button
           - Watch the progress bar as files are processed
           
        5. **Ask Questions**: 
           - Use the chat below to ask about your files
           - Mention specific filenames for detailed analysis
        """)
        
        st.markdown("### Supported File Types:")
        col1, col2, col3 = st.columns(3)
        
        with col1:
            st.markdown("""
            **Documents:**
            - PDF (.pdf)
            - Word (.docx)
            - PowerPoint (.pptx)
            - Excel (.xlsx, .xls)
            """)
        
        with col2:
            st.markdown("""
            **Code Files:**
            - Python (.py)
            - JavaScript (.js)
            - Java (.java)
            - HTML (.html)
            - CSS (.css)
            """)
        
        with col3:
            st.markdown("""
            **Data Files:**
            - Text (.txt)
            - Markdown (.md)
            - JSON (.json)
            - CSV (.csv)
            - XML (.xml)
            """)
    
    with st.expander("💡 Example Questions You Can Ask", expanded=False):
        st.markdown("""
        ### General Questions:
        - "What types of files did you find?"
        - "Give me an overview of all the documents"
        - "Which files contain code?"
        
        ### Specific File Questions:
        - "What does main.py do?" (if you have a main.py file)
        - "Summarize the contents of report.pdf"
        - "What are the main points in presentation.pptx?"
        - "Explain the data in spreadsheet.xlsx"
        
        ### Analysis Questions:
        - "Find any TODO comments in the code files"
        - "What are the main functions in my Python files?"
        - "Are there any important dates or deadlines mentioned?"
        """)

    # Logging example
    with st.expander("🔧 Logging Example", expanded=False):
        st.code("""
# Example of logging in action:
import logging

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# This application uses logging to track:
logger.info("Starting file scan...")
logger.warning("Large file detected - may use more API quota")
logger.error("Failed to read file: permission denied")
logger.debug("API response received successfully")
        """, language="python")
        
        if st.button("Test Logging Output"):
            # Create a container to show logs in the app
            log_container = st.container()
            
            with log_container:
                st.write("**Log Output:**")
                
                # Log messages with explicit flushing
                logger.info("✅ Example INFO log: Application started successfully")
                st.success("INFO: Application started successfully")
                
                logger.warning("⚠️ Example WARNING log: Processing large file")
                st.warning("WARNING: Processing large file")
                
                logger.error("❌ Example ERROR log: This is just a demonstration")
                st.error("ERROR: This is just a demonstration")
                
                # Force flush all handlers
                for handler in logger.handlers:
                    handler.flush()
                for handler in logging.getLogger().handlers:
                    handler.flush()
                
                # Show current working directory and log file location
                current_dir = os.getcwd()
                log_file_path = os.path.join(current_dir, 'app.log')
                
                st.info(f"Logs should be written to: `{log_file_path}`")
                st.info("Check your console/terminal for immediate log output!")
                
                # Check if log file exists and show its size
                if os.path.exists(log_file_path):
                    file_size = os.path.getsize(log_file_path)
                    st.success(f"✅ Log file exists! Size: {file_size} bytes")
                    
                    # Show last few lines of the log file
                    try:
                        with open(log_file_path, 'r', encoding='utf-8') as f:
                            lines = f.readlines()
                            if lines:
                                st.text_area("Last few log entries:", 
                                           value=''.join(lines[-10:]), 
                                           height=200)
                    except Exception as e:
                        st.error(f"Could not read log file: {e}")
                else:
                    st.error("❌ Log file not found!")

# --- Core AI and File Handling Functions ---

def make_api_call(api_key, prompt, model="gemini-2.0-flash"):
    """
    Generic function to make a call to the Gemini API.
    Handles the request and response logic.
    """
    if not api_key:
        return "[Error: API key not provided]"

    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={api_key}"
    
    payload = {
        "contents": [{"parts": [{"text": prompt}]}],
        "generationConfig": {
            "temperature": 0.5,
            "topP": 1.0,
            "maxOutputTokens": 2048,
        }
    }
    
    headers = {'Content-Type': 'application/json'}

    try:
        response = requests.post(url, headers=headers, json=payload, timeout=120)
        response.raise_for_status()
        result = response.json()
        
        if (candidates := result.get('candidates')) and (parts := candidates[0].get('content', {}).get('parts')):
            return parts[0].get('text', '[Could not extract text from response]').strip()
        else:
            return f"[API Error: Unexpected response format. Response: {result}]"

    except requests.exceptions.RequestException as e:
        return f"[API Request Error: {e}]"
    except (KeyError, IndexError):
        return "[API Error: Could not parse the summary from the API response]"


def get_file_content(file_path, full_content=False):
    """
    Extracts text content from various file types.
    Can be limited for initial summary or read fully for Q&A.
    """
    _, extension = os.path.splitext(file_path)
    extension = extension.lower()
    content = ""

    try:
        # Handling for plain text and code files
        if extension in ['.txt', '.md', '.py', '.java', '.js', '.html', '.css', '.json', '.xml', '.csv', '.log', '.ini', '.cfg', '.sh', '.bat']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
        elif extension == '.pdf':
            with fitz.open(file_path) as doc:
                content = "".join(page.get_text() for page in doc)
        elif extension == '.pptx':
            prs = pptx.Presentation(file_path)
            content = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        elif extension in ['.xlsx', '.xls']:
            workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            text_parts = []
            for sheet in workbook:
                for row in sheet.iter_rows():
                    text_parts.append(" ".join([str(cell.value) for cell in row if cell.value is not None]))
            content = "\n".join(text_parts)
        elif extension == '.docx':
            doc = docx.Document(file_path)
            content = "\n".join([para.text for para in doc.paragraphs])
        else:
            return None, "Unsupported file type"
    except Exception as e:
        return None, str(e)
    
    if full_content:
        return content, None
    else:
        # For initial summary, truncate to keep it free-tier friendly
        return content[:15000], None

# --- Streamlit UI and Application Logic ---

st.set_page_config(layout="wide", page_title="AI File Summarizer & Q&A")

st.title("📁 AI-Powered File Summarizer and Q&A")

# Initialize session state variables
if "messages" not in st.session_state:
    st.session_state.messages = []
if "scanned_files" not in st.session_state:
    st.session_state.scanned_files = {}

# --- Sidebar for Inputs ---
with st.sidebar:
    st.header("⚙️ Settings")
    api_key = st.text_input("Enter your Google Gemini API Key", type="password", help="Your API key is required.")
    path_input = st.text_input("Enter the full directory path", placeholder="e.g., C:\\Users\\YourUser\\Documents")

    if st.button("1. Scan & Summarize Directory", type="primary"):
        st.session_state.scanned_files = {}
        st.session_state.messages = []
        if not api_key:
            st.error("Please enter your API Key.")
        elif not path_input or not os.path.isdir(path_input):
            st.error("Please enter a valid directory path.")
        else:
            st.session_state.path_input = path_input
            st.success(f"Scanning: {os.path.abspath(path_input)}")
            # This just triggers a re-run; the actual scan happens in the main body
    
    st.markdown("---")
    st.markdown("After scanning, you can ask questions about the files in the chat below.")


# --- Main Content Area ---

# Perform the scan if the button was pressed
if "path_input" in st.session_state and not st.session_state.scanned_files:
    path = st.session_state.path_input
    supported_ext = ['.txt', '.md', '.py', '.java', '.js', '.html', '.css', '.json', 
                     '.xml', '.csv', '.log', '.ini', '.cfg', '.sh', '.bat',
                     '.pdf', '.pptx', '.xlsx', '.xls', '.docx']
    
    files_to_process = []
    for root, _, files in os.walk(path):
        for file_name in files:
            if os.path.splitext(file_name)[1].lower() in supported_ext:
                files_to_process.append((os.path.join(root, file_name), file_name))

    progress_bar = st.progress(0, text="Starting scan...")

    for i, (file_path, file_name) in enumerate(files_to_process):
        progress_text = f"Processing {i+1}/{len(files_to_process)}: {file_name}"
        progress_bar.progress((i + 1) / len(files_to_process), text=progress_text)
        
        # Get truncated content for initial summary
        content, error = get_file_content(file_path)
        summary = "[Could not read file]"
        
        # FIX: Check if content is valid before making an API call
        if not error and content and not content.isspace():
            summary_prompt = f"Provide a concise, one-paragraph summary for a file named '{file_name}'. Focus on its main purpose and key topics.\n\n---\n\n{content}"
            summary = make_api_call(api_key, summary_prompt)
        elif not error:
             summary = "[File is empty or contains no readable text]"


        st.session_state.scanned_files[file_name] = {
            "path": file_path,
            "summary": summary
        }
    
    progress_bar.empty()
    st.success("Scan complete! You can now ask questions below.")


# Display initial summaries if available
if st.session_state.scanned_files:
    st.markdown("---")
    st.header("Initial Summaries")
    for file_name, data in st.session_state.scanned_files.items():
        with st.expander(f"📄 **{file_name}**"):
            st.markdown(f"> {data['summary']}")

    st.markdown("---")
    st.header("💬 Chat About Your Files")
    
    # Display chat history
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])

    # Chat input
    if prompt := st.chat_input("Ask a question about your files..."):
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)

        with st.chat_message("assistant"):
            message_placeholder = st.empty()
            with st.spinner("Thinking..."):
                # Check if the prompt mentions a specific file
                mentioned_file = None
                for file_name in st.session_state.scanned_files.keys():
                    if re.search(r'\b' + re.escape(file_name) + r'\b', prompt, re.IGNORECASE):
                        mentioned_file = file_name
                        break
                
                final_prompt = prompt
                if mentioned_file:
                    file_path = st.session_state.scanned_files[mentioned_file]["path"]
                    
                    # --- COST WARNING ---
                    message_placeholder.warning(f"⚠️ You mentioned **{mentioned_file}**. "
                                                f"I am now reading the **ENTIRE FILE** to give a detailed answer. "
                                                f"If this is a large file, this may use a significant portion of your free API quota and could incur costs.")
                    
                    full_content, error = get_file_content(file_path, full_content=True)
                    if error:
                        response = f"I'm sorry, I couldn't read the full content of {mentioned_file}. Error: {error}"
                    else:
                        final_prompt = (f"You are a helpful assistant. The user is asking a question about a specific file. "
                                        f"Here is the user's question: '{prompt}'\n\n"
                                        f"And here is the FULL text content of the file named '{mentioned_file}':\n\n---\n\n"
                                        f"{full_content}\n\n---\n\n"
                                        f"Please answer the user's question based on the provided file content.")
                        response = make_api_call(api_key, final_prompt)
                else:
                    # General question, not about a specific file
                    response = make_api_call(api_key, prompt)
                
                message_placeholder.markdown(response)
        
        st.session_state.messages.append({"role": "assistant", "content": response})

else:
    st.info("Click the 'Scan & Summarize Directory' button in the sidebar to begin.")
